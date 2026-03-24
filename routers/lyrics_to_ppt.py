from fastapi import APIRouter
from fastapi.responses import Response
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import lyricsgenius
import re
import math
import io
from urllib.parse import quote

router = APIRouter()

# Tvůj původní token
MY_GENIUS_TOKEN = "YCvV3ZCX6orJkY9A4Jl6k6dl-fS55dbQTKhQ4XUIzjSqTwKy7Cbk69qhhxrKCkWv"

# --- Šablony pro to, co nám pošle Frontend ---
class SearchRequest(BaseModel):
    interpret: str
    pisen: str

class PPTXRequest(BaseModel):
    text_pisne: str
    max_radku: int = 4
    nazev_souboru: str = "Píseň"

# --- 1. API PRO VYHLEDÁNÍ TEXTU ---
@router.post("/api/vyhledat-text")
def search_lyrics(req: SearchRequest):
    try:
        genius = lyricsgenius.Genius(MY_GENIUS_TOKEN)
        song_data = genius.search_song(req.pisen, req.interpret)

        if song_data:
            # Tvoje původní čistící logika
            text = re.sub(r"^\d+ Contributors", "", song_data.lyrics)
            text = re.sub(r"\[.*?\]", "", text)
            text = re.sub(r"\d*Embed$", "", text)
            text = text.replace('\r\n', '\n')
            text = re.sub(r'\n\s*\n', '\n\n', text)
            return {"text": text.strip()}
        else:
            return {"chyba": "Píseň nebyla nalezena."}
    except Exception as e:
        return {"chyba": str(e)}

# --- 2. API PRO VYTVOŘENÍ PREZENTACE ---
@router.post("/api/vytvorit-prezentaci")
def create_pptx(req: PPTXRequest):
    blocks = req.text_pisne.split("\n\n")
    blocks = [b.strip() for b in blocks if b.strip()]

    if not blocks:
        return {"chyba": "Žádný text k převodu."}

    # Vytvoření prezentace
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide(text_content):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(6.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text_content
        p.font.size = Pt(44)
        p.font.name = 'Arial'
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    for block in blocks:
        lines = block.splitlines()
        if len(lines) <= req.max_radku:
            add_slide("\n".join(lines))
        else:
            num_parts = math.ceil(len(lines) / req.max_radku)
            chunk_size = math.ceil(len(lines) / num_parts)
            for i in range(0, len(lines), chunk_size):
                chunk = lines[i: i + chunk_size]
                add_slide("\n".join(chunk))

    # Uložení do virtuální paměti a odeslání prohlížeči
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    
    bezpecny_nazev = quote(req.nazev_souboru + ".pptx")

    return Response(
        content=ppt_stream.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{bezpecny_nazev}"}
    )
